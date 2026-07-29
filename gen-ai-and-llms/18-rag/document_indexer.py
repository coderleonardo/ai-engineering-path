"""Multi-format document indexing for RAG — Reference Script.

Reference script. See `rag.md` for the ingestion/query pipeline this implements and the chunking
strategy behind the parameters below, and `rag_api.py`/`rag_streamlit_app.py` for the retrieval side
that queries the collection this script builds.

Methods covered:
- Recursively walking a directory and extracting text per file type (PDF via `pypdf`, `.txt`, `.docx`
  via `python-docx`, `.pptx` via `python-pptx`)
- Token-based chunking (`TokenTextSplitter`)
- Embedding chunks locally (`HuggingFaceEmbeddings`) and upserting them into a Qdrant collection
  (`QdrantVectorStore`), with per-chunk metadata that later enables citing the source file

Use this as a reference when: you need a from-scratch multi-format ingestion pipeline for a RAG
system — dispatch-by-extension text extraction, chunking, embedding, and vector store upsert, without
a heavier document-loading framework.

Don't use this as a reference for: the query/retrieval side (see `rag_api.py`), or production-grade
document parsing (this uses lightweight per-format libraries, not something like `unstructured` that
handles tables, images, and layout).

Requires (not part of the repo's shared pyproject.toml — install separately): pypdf, python-docx,
python-pptx, langchain-text-splitters, langchain-huggingface, sentence-transformers, qdrant-client,
langchain-qdrant.

Run: python document_indexer.py <path-to-documents-directory>
Requires a Qdrant instance reachable at QDRANT_URL (default http://localhost:6333), e.g.:
  docker run -p 6333:6333 qdrant/qdrant
"""

import os
import sys
from os import listdir
from os.path import isdir, isfile, join

import docx
from pptx import Presentation
from pypdf import PdfReader
from langchain_text_splitters import TokenTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

import warnings
warnings.filterwarnings("ignore")

QDRANT_URL = os.environ.get("QDRANT_URL", "http://localhost:6333")
COLLECTION_NAME = "documents"

# This embedding model is trained with a dot-product ranking objective (not cosine), so the Qdrant
# collection below is created with Distance.DOT to match -- pairing a different embedding model with
# this distance metric (or vice versa) would silently degrade retrieval quality.
EMBEDDING_MODEL = "sentence-transformers/msmarco-bert-base-dot-v5"
EMBEDDING_DIM = 768


def list_files(directory):
    """Returns every file path under `directory`, recursing into subdirectories."""
    paths = []
    for entry in listdir(directory):
        full_path = join(directory, entry)
        if isfile(full_path):
            paths.append(full_path)
        elif isdir(full_path):
            paths += list_files(full_path)
    return paths


def load_docx_text(path):
    document = docx.Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def load_pptx_text(path):
    presentation = Presentation(path)
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return "\n".join(lines)


def load_pdf_text(path):
    reader = PdfReader(path)
    return " ".join(page.extract_text() for page in reader.pages)


# Dispatch by extension -- each loader returns the file's full text as a single string, chunked below.
LOADERS = {
    ".pdf": load_pdf_text,
    ".txt": lambda path: open(path, "r").read(),
    ".docx": load_docx_text,
    ".pptx": load_pptx_text,
}


def index_documents(directory):
    embeddings = HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )

    client = QdrantClient(QDRANT_URL)

    # Recreated from scratch on every run -- this is a full-corpus reindex, not an incremental update.
    # An incremental pipeline would diff against what's already stored instead of dropping the collection.
    if client.collection_exists(COLLECTION_NAME):
        client.delete_collection(COLLECTION_NAME)
    client.create_collection(
        COLLECTION_NAME,
        vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.DOT),
    )

    vector_store = QdrantVectorStore(client=client, collection_name=COLLECTION_NAME, embedding=embeddings)

    # chunk_size/chunk_overlap are in tokens (TokenTextSplitter), not characters -- this keeps each
    # chunk within a predictable token budget for the LLM context window regardless of how verbose the
    # source text is. The overlap keeps a fact that lands near a chunk boundary from being fully lost.
    text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)

    print("\nIndexing documents...\n")

    for path in list_files(directory):
        extension = "." + path.rsplit(".", 1)[-1].lower() if "." in path else ""
        loader = LOADERS.get(extension)
        if loader is None:
            continue

        try:
            print("Indexing:", path)
            content = loader(path)
            chunks = text_splitter.split_text(content)

            # Attaching the source path to every chunk's metadata is what lets the API cite "which
            # document did this answer come from" later -- see rag_api.py.
            metadatas = [{"path": path} for _ in chunks]
            vector_store.add_texts(chunks, metadatas=metadatas)
        except Exception as e:
            print(f"Failed to index {path}: {e}")

    print("\nIndexing complete.\n")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        index_documents(sys.argv[1])
    else:
        print("Usage: python document_indexer.py <path-to-documents-directory>")
